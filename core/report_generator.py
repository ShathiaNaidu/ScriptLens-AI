from __future__ import annotations

import json
from html import escape
from io import BytesIO
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

from reportlab.platypus import (
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from core.models import AnalysisReport


APP_TARGET_MARKET = [
    "Student filmmakers and final-year project teams",
    "Independent screenwriters",
    "Film and media students",
    "Lecturers, film schools, and universities",
    "Independent producers and small production companies",
    "Film competitions and script-development programmes",
]


def _safe_filename(value: str) -> str:
    allowed = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_"
    cleaned = "".join(ch if ch in allowed else "_" for ch in value).strip("_")
    return cleaned or "screenplay"


def report_basename(report: AnalysisReport) -> str:
    return f"{_safe_filename(report.metadata.title)}_Screenplay_Report"


def pitch_deck_basename(report: AnalysisReport) -> str:
    return f"{_safe_filename(report.metadata.title)}_Investor_Pitch_Deck"


def generate_json(report: AnalysisReport) -> bytes:
    return json.dumps(
        report.model_dump(mode="json"), ensure_ascii=False, indent=2
    ).encode("utf-8")


def _add_bullets_docx(document: Document, items: Iterable[str]) -> None:
    for item in items:
        document.add_paragraph(str(item), style="List Bullet")


def _set_docx_defaults(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10)
    for style_name in ("Title", "Heading 1", "Heading 2", "Heading 3"):
        styles[style_name].font.name = "Aptos Display"


def generate_docx(report: AnalysisReport) -> bytes:
    document = Document()
    _set_docx_defaults(document)

    title = document.add_heading(f"{report.metadata.title} - Screenplay Analysis", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = document.add_paragraph(
        f"Writer: {report.metadata.writer} | Overall score: {report.scores.overall_score}/100"
    )
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_heading("Executive Summary", level=1)
    document.add_paragraph(report.main_recommendation)
    document.add_paragraph(f"Logline: {report.metadata.logline}")
    document.add_paragraph(f"Story concept: {report.metadata.story_concept}")

    document.add_heading("Screenplay Overview", level=1)
    overview = document.add_table(rows=0, cols=2)
    overview.style = "Table Grid"
    overview_data = [
        ("Title", report.metadata.title),
        ("Writer", report.metadata.writer),
        ("Genre", " / ".join(report.metadata.genres)),
        ("Format", report.metadata.format),
        ("Estimated runtime", f"{report.metadata.estimated_runtime_minutes} minutes"),
        ("Pages", str(report.metadata.page_count)),
        ("Scenes", str(report.metadata.scene_count)),
        ("Languages", ", ".join(report.metadata.languages_detected)),
        ("Central conflict", report.metadata.central_conflict),
    ]
    for key, value in overview_data:
        cells = overview.add_row().cells
        cells[0].text = key
        cells[1].text = value or "-"

    document.add_heading("Scores", level=1)
    scores = report.scores.model_dump()
    score_table = document.add_table(rows=1, cols=2)
    score_table.style = "Table Grid"
    score_table.rows[0].cells[0].text = "Category"
    score_table.rows[0].cells[1].text = "Score"
    for category, score in scores.items():
        cells = score_table.add_row().cells
        cells[0].text = category.replace("_", " ").title()
        cells[1].text = f"{score}/100"

    document.add_heading("Top Strengths", level=1)
    _add_bullets_docx(document, report.top_strengths)
    document.add_heading("Priority Improvements", level=1)
    _add_bullets_docx(document, report.priority_improvements)

    document.add_heading("Character Analysis", level=1)
    for character in report.characters:
        document.add_heading(character.name, level=2)
        document.add_paragraph(f"Role: {character.role}")
        document.add_paragraph(f"Traits: {', '.join(character.traits)}")
        document.add_paragraph(f"Goal: {character.goal}")
        document.add_paragraph(f"Conflict: {character.conflict}")
        document.add_paragraph(
            f"Arc: {character.arc.beginning} -> {character.arc.middle} -> {character.arc.ending}"
        )
        document.add_paragraph(f"AI feedback: {character.ai_feedback}")
        if character.strengths:
            document.add_paragraph("Strengths:")
            _add_bullets_docx(document, character.strengths)
        if character.improvements:
            document.add_paragraph("Improvements:")
            _add_bullets_docx(document, character.improvements)

    document.add_heading("Story Structure", level=1)
    for act in report.acts:
        document.add_heading(f"Act {act.act_number}: {act.title}", level=2)
        document.add_paragraph(f"Scene range: {act.approximate_scene_range}")
        document.add_paragraph(f"Purpose: {act.purpose}")
        document.add_paragraph("Key events:")
        _add_bullets_docx(document, act.key_events)
        document.add_paragraph("Suggestions:")
        _add_bullets_docx(document, act.suggestions)

    document.add_heading("Scene Analysis", level=1)
    scene_table = document.add_table(rows=1, cols=6)
    scene_table.style = "Table Grid"
    headers = ["Scene", "Heading", "Purpose", "Emotion", "Suspense", "Main suggestion"]
    for index, header in enumerate(headers):
        scene_table.rows[0].cells[index].text = header
    for scene in report.scenes:
        cells = scene_table.add_row().cells
        cells[0].text = str(scene.scene_number)
        cells[1].text = scene.heading
        cells[2].text = scene.purpose
        cells[3].text = scene.dominant_emotion
        cells[4].text = f"{scene.suspense_score}/10"
        cells[5].text = scene.suggestions[0] if scene.suggestions else "-"

    document.add_heading("Dialogue Analysis", level=1)
    for dialogue in report.dialogue_analysis:
        heading = f"{dialogue.speaker} - Scene {dialogue.scene_number or 'Unknown'}"
        document.add_heading(heading, level=2)
        document.add_paragraph(f'"{dialogue.dialogue_excerpt}"')
        document.add_paragraph(f"Purpose: {dialogue.purpose}")
        document.add_paragraph(f"Voice match: {dialogue.character_voice_match}")
        document.add_paragraph(f"Naturalness: {dialogue.naturalness_score}/10")
        _add_bullets_docx(document, dialogue.improvements)

    document.add_heading("Genre Analysis", level=1)
    for genre in report.genre_analysis:
        document.add_heading(f"{genre.genre}: {genre.score}/100", level=2)
        document.add_paragraph(genre.reason)
        _add_bullets_docx(document, genre.suggestions)

    document.add_heading("Originality and Local Identity", level=1)
    document.add_paragraph(report.originality.originality_summary)
    document.add_paragraph("Familiar storytelling patterns:")
    _add_bullets_docx(document, report.originality.familiar_storytelling_patterns)
    document.add_paragraph("Distinctive elements:")
    _add_bullets_docx(document, report.originality.distinctive_elements)
    document.add_paragraph("Local identity opportunities:")
    _add_bullets_docx(document, report.originality.local_identity_opportunities)
    document.add_paragraph(f"Important: {report.originality.disclaimer}")

    document.add_heading("Audience Prediction", level=1)
    audience_table = document.add_table(rows=1, cols=4)
    audience_table.style = "Table Grid"
    for index, header in enumerate(["Audience", "Appeal", "Score", "Reason"]):
        audience_table.rows[0].cells[index].text = header
    for segment in report.audience_prediction:
        cells = audience_table.add_row().cells
        cells[0].text = segment.segment
        cells[1].text = segment.predicted_appeal
        cells[2].text = f"{segment.appeal_score}/100"
        cells[3].text = segment.reason

    document.add_heading("Step 13 - Storyboard Generator", level=1)
    for panel in report.storyboard:
        document.add_heading(f"Scene {panel.scene_number}: {panel.title}", level=2)
        document.add_paragraph(f"Visual: {panel.visual_description}")
        document.add_paragraph(f"Camera angle: {panel.camera_angle}")
        document.add_paragraph(f"Shot type: {panel.shot_type}")
        document.add_paragraph(f"Character positions: {panel.character_positions}")
        document.add_paragraph(f"Lighting: {panel.lighting}")
        document.add_paragraph(f"Mood: {panel.mood}")
        document.add_paragraph(f"AI concept-art prompt: {panel.concept_art_prompt}")

    document.add_heading("Step 14 - Pitch Generator", level=1)
    pitch_package = report.pitch_package
    document.add_paragraph(f"Logline: {pitch_package.logline}")
    document.add_paragraph(f"One-page synopsis: {pitch_package.one_page_synopsis}")
    document.add_heading("Character Profiles", level=2)
    for character in pitch_package.character_profiles:
        document.add_paragraph(f"{character.name} - {character.role}: {character.pitch_description}")
    document.add_heading("Director's Vision", level=2)
    document.add_paragraph(pitch_package.directors_vision)
    document.add_heading("Mood Board", level=2)
    _add_bullets_docx(document, pitch_package.mood_board)
    document.add_paragraph(f"Budget estimate: {pitch_package.budget_estimate}")
    document.add_heading("Target Audience", level=2)
    _add_bullets_docx(document, pitch_package.target_audience)
    document.add_heading("Suggested Platforms", level=2)
    _add_bullets_docx(document, pitch_package.suggested_platforms)
    document.add_heading("Marketing Strategy", level=2)
    _add_bullets_docx(document, pitch_package.marketing_strategy)
    document.add_heading("Poster Concept", level=2)
    document.add_paragraph(pitch_package.poster_concept)
    document.add_paragraph(f"Poster art prompt: {pitch_package.poster_art_prompt}")
    document.add_heading("Investor Pitch Deck", level=2)
    for slide in pitch_package.investor_pitch_deck:
        document.add_paragraph(f"Slide {slide.slide_number}: {slide.title}")
        _add_bullets_docx(document, slide.key_points)
    document.add_heading("Final AI Report", level=2)
    final_scores = pitch_package.final_scores.model_dump()
    for category, score in final_scores.items():
        document.add_paragraph(f"{category.replace('_', ' ').title()}: {score}/100")
    document.add_paragraph(f"AI Recommendation: {pitch_package.ai_recommendation}")

    document.add_heading("Producer Pitch", level=1)
    document.add_paragraph(f"Logline: {report.producer_pitch.logline}")
    document.add_paragraph(report.producer_pitch.pitch_paragraph)
    document.add_paragraph(f"Short synopsis: {report.producer_pitch.short_synopsis}")
    document.add_paragraph(f"Genre: {report.producer_pitch.genre}")
    document.add_paragraph("Target audience:")
    _add_bullets_docx(document, report.producer_pitch.target_audience)
    document.add_paragraph("Target market:")
    _add_bullets_docx(document, report.producer_pitch.target_market)
    document.add_paragraph("Selling points:")
    _add_bullets_docx(document, report.producer_pitch.selling_points)
    document.add_paragraph("Production considerations:")
    _add_bullets_docx(document, report.producer_pitch.production_considerations)

    document.add_heading("Target Market of the Cinevora App", level=1)
    _add_bullets_docx(document, APP_TARGET_MARKET)

    document.add_heading("Analysis Limitations", level=1)
    _add_bullets_docx(document, report.analysis_limitations)

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _register_pdf_font() -> str:
    candidates = [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
    ]
    for path in candidates:
        if path.exists():
            try:
                pdfmetrics.registerFont(TTFont("CinevoraFont", str(path)))
                return "CinevoraFont"
            except Exception:
                continue
    return "Helvetica"


def _pdf_bullets(items: Iterable[str], style: ParagraphStyle) -> list:
    flowables: list = []
    for item in items:
        flowables.append(Paragraph(f"&#8226; {escape(str(item))}", style))
    return flowables


def generate_pdf(report: AnalysisReport) -> bytes:
    buffer = BytesIO()
    font_name = _register_pdf_font()
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="SLTitle",
            parent=styles["Title"],
            fontName=font_name,
            fontSize=22,
            leading=27,
            alignment=TA_CENTER,
            spaceAfter=12,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SLHeading1",
            parent=styles["Heading1"],
            fontName=font_name,
            fontSize=15,
            leading=19,
            spaceBefore=12,
            spaceAfter=7,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SLHeading2",
            parent=styles["Heading2"],
            fontName=font_name,
            fontSize=12,
            leading=15,
            spaceBefore=8,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SLBody",
            parent=styles["BodyText"],
            fontName=font_name,
            fontSize=9.2,
            leading=13,
            spaceAfter=5,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SLSmall",
            parent=styles["BodyText"],
            fontName=font_name,
            fontSize=8,
            leading=10,
        )
    )

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=16 * mm,
        rightMargin=16 * mm,
        topMargin=15 * mm,
        bottomMargin=15 * mm,
        title=f"{report.metadata.title} Screenplay Report",
        author="Cinevora AI",
    )
    story: list = []
    story.append(Paragraph(escape(report.metadata.title), styles["SLTitle"]))
    story.append(
        Paragraph(
            escape(
                f"Screenplay Analysis | Writer: {report.metadata.writer} | Overall: {report.scores.overall_score}/100"
            ),
            styles["SLBody"],
        )
    )
    story.append(Spacer(1, 6))
    story.append(Paragraph("Executive Summary", styles["SLHeading1"]))
    story.append(Paragraph(escape(report.main_recommendation), styles["SLBody"]))
    story.append(Paragraph(f"<b>Logline:</b> {escape(report.metadata.logline)}", styles["SLBody"]))

    story.append(Paragraph("Screenplay Overview", styles["SLHeading1"]))
    overview_data = [
        ["Title", report.metadata.title],
        ["Writer", report.metadata.writer],
        ["Genre", " / ".join(report.metadata.genres)],
        ["Format", report.metadata.format],
        ["Runtime", f"{report.metadata.estimated_runtime_minutes} minutes"],
        ["Pages / scenes", f"{report.metadata.page_count} / {report.metadata.scene_count}"],
        ["Locations", ", ".join(report.metadata.locations)],
        ["Central conflict", report.metadata.central_conflict],
    ]
    overview_table = Table(
        [[Paragraph(escape(str(c)), styles["SLSmall"]) for c in row] for row in overview_data],
        colWidths=[42 * mm, 130 * mm],
        repeatRows=0,
    )
    overview_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#EEF2FF")),
                ("FONTNAME", (0, 0), (-1, -1), font_name),
                ("LEFTPADDING", (0, 0), (-1, -1), 5),
                ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    story.append(overview_table)

    story.append(Paragraph("Scores", styles["SLHeading1"]))
    score_rows = [["Category", "Score"]]
    for category, score in report.scores.model_dump().items():
        score_rows.append([category.replace("_", " ").title(), f"{score}/100"])
    score_table = Table(score_rows, colWidths=[120 * mm, 35 * mm], repeatRows=1)
    score_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, -1), font_name),
                ("ALIGN", (1, 1), (1, -1), "CENTER"),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    story.append(score_table)

    story.append(Paragraph("Top Strengths", styles["SLHeading1"]))
    story.extend(_pdf_bullets(report.top_strengths, styles["SLBody"]))
    story.append(Paragraph("Priority Improvements", styles["SLHeading1"]))
    story.extend(_pdf_bullets(report.priority_improvements, styles["SLBody"]))

    story.append(PageBreak())
    story.append(Paragraph("Character Analysis", styles["SLHeading1"]))
    for character in report.characters:
        block = [
            Paragraph(escape(character.name), styles["SLHeading2"]),
            Paragraph(f"<b>Role:</b> {escape(character.role)}", styles["SLBody"]),
            Paragraph(f"<b>Traits:</b> {escape(', '.join(character.traits))}", styles["SLBody"]),
            Paragraph(f"<b>Goal:</b> {escape(character.goal)}", styles["SLBody"]),
            Paragraph(f"<b>Conflict:</b> {escape(character.conflict)}", styles["SLBody"]),
            Paragraph(
                f"<b>Arc:</b> {escape(character.arc.beginning)} -&gt; {escape(character.arc.middle)} -&gt; {escape(character.arc.ending)}",
                styles["SLBody"],
            ),
            Paragraph(f"<b>Feedback:</b> {escape(character.ai_feedback)}", styles["SLBody"]),
        ]
        story.append(KeepTogether(block))

    story.append(Paragraph("Story Structure", styles["SLHeading1"]))
    for act in report.acts:
        story.append(Paragraph(f"Act {act.act_number}: {escape(act.title)}", styles["SLHeading2"]))
        story.append(Paragraph(f"<b>Purpose:</b> {escape(act.purpose)}", styles["SLBody"]))
        story.extend(_pdf_bullets(act.key_events, styles["SLBody"]))
        story.extend(_pdf_bullets(act.suggestions, styles["SLBody"]))

    story.append(PageBreak())
    story.append(Paragraph("Scene Analysis", styles["SLHeading1"]))
    scene_rows = [["Scene", "Heading", "Purpose", "Emotion", "Suspense", "Suggestion"]]
    for scene in report.scenes:
        scene_rows.append(
            [
                str(scene.scene_number),
                scene.heading,
                scene.purpose,
                scene.dominant_emotion,
                f"{scene.suspense_score}/10",
                scene.suggestions[0] if scene.suggestions else "-",
            ]
        )
    wrapped_scene_rows = [
        [Paragraph(escape(str(value)), styles["SLSmall"]) for value in row]
        for row in scene_rows
    ]
    scene_table = Table(
        wrapped_scene_rows,
        colWidths=[12 * mm, 38 * mm, 34 * mm, 24 * mm, 18 * mm, 48 * mm],
        repeatRows=1,
    )
    scene_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTNAME", (0, 0), (-1, -1), font_name),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    story.append(scene_table)

    story.append(Paragraph("Genre Analysis", styles["SLHeading1"]))
    for genre in report.genre_analysis:
        story.append(Paragraph(f"{escape(genre.genre)}: {genre.score}/100", styles["SLHeading2"]))
        story.append(Paragraph(escape(genre.reason), styles["SLBody"]))
        story.extend(_pdf_bullets(genre.suggestions, styles["SLBody"]))

    story.append(Paragraph("Originality and Audience", styles["SLHeading1"]))
    story.append(Paragraph(escape(report.originality.originality_summary), styles["SLBody"]))
    story.extend(_pdf_bullets(report.originality.distinctive_elements, styles["SLBody"]))
    story.append(Paragraph(f"<b>Disclaimer:</b> {escape(report.originality.disclaimer)}", styles["SLBody"]))
    for segment in report.audience_prediction:
        story.append(
            Paragraph(
                f"<b>{escape(segment.segment)} ({segment.appeal_score}/100):</b> {escape(segment.reason)}",
                styles["SLBody"],
            )
        )

    story.append(PageBreak())
    story.append(Paragraph("Step 13 - Storyboard Generator", styles["SLHeading1"]))
    for panel in report.storyboard:
        story.append(Paragraph(f"Scene {panel.scene_number}: {escape(panel.title)}", styles["SLHeading2"]))
        story.append(Paragraph(f"<b>Visual:</b> {escape(panel.visual_description)}", styles["SLBody"]))
        story.append(Paragraph(f"<b>Camera:</b> {escape(panel.camera_angle)} | <b>Shot:</b> {escape(panel.shot_type)}", styles["SLBody"]))
        story.append(Paragraph(f"<b>Blocking:</b> {escape(panel.character_positions)}", styles["SLBody"]))
        story.append(Paragraph(f"<b>Lighting:</b> {escape(panel.lighting)} | <b>Mood:</b> {escape(panel.mood)}", styles["SLBody"]))
        story.append(Paragraph(f"<b>Concept-art prompt:</b> {escape(panel.concept_art_prompt)}", styles["SLSmall"]))

    pitch_package = report.pitch_package
    story.append(PageBreak())
    story.append(Paragraph("Step 14 - Pitch Generator", styles["SLHeading1"]))
    story.append(Paragraph(f"<b>Logline:</b> {escape(pitch_package.logline)}", styles["SLBody"]))
    story.append(Paragraph(f"<b>One-page synopsis:</b> {escape(pitch_package.one_page_synopsis)}", styles["SLBody"]))
    story.append(Paragraph("Character Profiles", styles["SLHeading2"]))
    for character in pitch_package.character_profiles:
        story.append(Paragraph(f"<b>{escape(character.name)} - {escape(character.role)}:</b> {escape(character.pitch_description)}", styles["SLBody"]))
    story.append(Paragraph("Director's Vision", styles["SLHeading2"]))
    story.append(Paragraph(escape(pitch_package.directors_vision), styles["SLBody"]))
    story.append(Paragraph("Mood Board", styles["SLHeading2"]))
    story.extend(_pdf_bullets(pitch_package.mood_board, styles["SLBody"]))
    story.append(Paragraph(f"<b>Budget estimate:</b> {escape(pitch_package.budget_estimate)}", styles["SLBody"]))
    story.append(Paragraph("Target Audience", styles["SLHeading2"]))
    story.extend(_pdf_bullets(pitch_package.target_audience, styles["SLBody"]))
    story.append(Paragraph("Suggested Platforms", styles["SLHeading2"]))
    story.extend(_pdf_bullets(pitch_package.suggested_platforms, styles["SLBody"]))
    story.append(Paragraph("Marketing Strategy", styles["SLHeading2"]))
    story.extend(_pdf_bullets(pitch_package.marketing_strategy, styles["SLBody"]))
    story.append(Paragraph("Poster Concept", styles["SLHeading2"]))
    story.append(Paragraph(escape(pitch_package.poster_concept), styles["SLBody"]))
    story.append(Paragraph(f"<b>Poster art prompt:</b> {escape(pitch_package.poster_art_prompt)}", styles["SLSmall"]))
    story.append(Paragraph("Investor Pitch Deck", styles["SLHeading2"]))
    for slide in pitch_package.investor_pitch_deck:
        story.append(Paragraph(f"Slide {slide.slide_number}: {escape(slide.title)}", styles["SLBody"]))
        story.extend(_pdf_bullets(slide.key_points, styles["SLSmall"]))
    story.append(Paragraph("Final AI Report", styles["SLHeading2"]))
    for category, score in pitch_package.final_scores.model_dump().items():
        story.append(Paragraph(f"{escape(category.replace('_', ' ').title())}: {score}/100", styles["SLBody"]))
    story.append(Paragraph(f"<b>AI Recommendation:</b> {escape(pitch_package.ai_recommendation)}", styles["SLBody"]))

    story.append(PageBreak())
    story.append(Paragraph("Producer Pitch", styles["SLHeading1"]))
    story.append(Paragraph(f"<b>Logline:</b> {escape(report.producer_pitch.logline)}", styles["SLBody"]))
    story.append(Paragraph(escape(report.producer_pitch.pitch_paragraph), styles["SLBody"]))
    story.append(Paragraph(f"<b>Synopsis:</b> {escape(report.producer_pitch.short_synopsis)}", styles["SLBody"]))
    story.append(Paragraph("Selling Points", styles["SLHeading2"]))
    story.extend(_pdf_bullets(report.producer_pitch.selling_points, styles["SLBody"]))
    story.append(Paragraph("Production Considerations", styles["SLHeading2"]))
    story.extend(_pdf_bullets(report.producer_pitch.production_considerations, styles["SLBody"]))

    story.append(Paragraph("Target Market of the Cinevora App", styles["SLHeading1"]))
    story.extend(_pdf_bullets(APP_TARGET_MARKET, styles["SLBody"]))
    story.append(Paragraph("Analysis Limitations", styles["SLHeading1"]))
    story.extend(_pdf_bullets(report.analysis_limitations, styles["SLBody"]))

    doc.build(story)
    return buffer.getvalue()


def generate_pitch_pptx(report: AnalysisReport) -> bytes:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333333)
    presentation.slide_height = PptxInches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = report.metadata.title
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Investor Pitch Deck | {report.pitch_package.logline}"

    for item in report.pitch_package.investor_pitch_deck:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = item.title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, point in enumerate(item.key_points):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = point
            paragraph.font.size = PptxPt(22)
            paragraph.space_after = PptxPt(8)

    recommendation_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    recommendation_slide.shapes.title.text = "AI Recommendation"
    tf = recommendation_slide.placeholders[1].text_frame
    tf.text = report.pitch_package.ai_recommendation
    tf.paragraphs[0].font.size = PptxPt(22)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
